"""
PowerPoint text extraction services for both .pptx and .ppt formats
Supports modern .pptx using python-pptx and legacy .ppt using Java converter
"""
import os
import subprocess
import tempfile
from typing import List, Dict, Any
from .structured_logger import logger, OperationTimer
from .file_validator import validate_file
from io import BytesIO
from pptx import Presentation


def extract_from_shape(shape, collected):
    """Extract text from PowerPoint shapes recursively"""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            collected.append(paragraph.text)

    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    collected.append(paragraph.text)

    if shape.shape_type == 6:  # GROUPED SHAPE
        for subshape in shape.shapes:
            extract_from_shape(subshape, collected)


def extract_text_from_pptx_file(file_like):
    """
    Extract text from modern PowerPoint files (.pptx)
    
    Args:
        file_like: File-like object containing PowerPoint data
        
    Returns:
        List of dictionaries with slide number and text content
    """
    prs = Presentation(file_like)
    all_text = []

    for slide_index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            extract_from_shape(shape, slide_text)

        combined = "\n".join(filter(None, slide_text))
        all_text.append({"slide": slide_index + 1, "text": combined})
    
    return all_text


def extract_text_from_ppt_file(file_path):
    """
    Extract text from legacy PowerPoint files (.ppt) using Java converter
    
    Args:
        file_path: Path to the .ppt file
        
    Returns:
        List of dictionaries with slide number and text content
    """
    jar_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "ppt_converter/target/ppt-converter-1.0-jar-with-dependencies.jar")

    try:
        result = subprocess.run(
            ["java", "-jar", jar_path, file_path],
            capture_output=True,
            text=True,
            check=True
        )
    except subprocess.CalledProcessError as e:
        logger.error(
            "Java PowerPoint extraction failed",
            file_path=file_path,
            returncode=e.returncode,
            stdout=e.stdout,
            stderr=e.stderr
        )
        return []

    slides_text = []
    current_slide = 0
    slide_lines = []

    for line in result.stdout.splitlines():
        if line.startswith("--- Slide"):
            if slide_lines:
                slides_text.append({"slide": current_slide, "text": "\n".join(slide_lines)})
                slide_lines = []

            current_slide = int(line.replace("--- Slide", "").replace("---", "").strip())
        else:
            slide_lines.append(line)

    if slide_lines:
        slides_text.append({"slide": current_slide, "text": "\n".join(slide_lines)})

    return slides_text


def extract_powerpoint_text(file_bytes: bytes, filename: str, request_id: str = None):
    """
    Universal PowerPoint text extraction
    
    Args:
        file_bytes: PowerPoint file content as bytes
        filename: Name of the file
        request_id: Optional request ID for tracing
        
    Returns:
        List of dictionaries with slide number and text content
    """
    with OperationTimer("powerpoint_extraction", filename=filename, request_id=request_id):
        try:
            # Validate file
            validation_result = validate_file(
                filename or "",
                len(file_bytes),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
            if not validation_result["valid"]:
                logger.log_validation_error(
                    filename=filename or "unknown",
                    errors=validation_result["errors"],
                    request_id=request_id
                )
                raise ValueError(f"File validation failed: {validation_result['errors']}")
            
            logger.info(
                "Starting PowerPoint extraction",
                filename=filename,
                file_size=len(file_bytes),
                request_id=request_id
            )
            
            filename = filename.lower()
            
            if filename.endswith(".pptx"):
                file_like = BytesIO(file_bytes)
                result = extract_text_from_pptx_file(file_like)
                
            elif filename.endswith(".ppt"):
                # Save to a secure temp file for Java extractor
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".ppt")
                try:
                    tmp.write(file_bytes)
                    tmp.flush()
                    tmp.close()
                    result = extract_text_from_ppt_file(tmp.name)
                finally:
                    try:
                        os.remove(tmp.name)
                    except Exception:
                        logger.warning(
                            "Failed to remove temp .ppt file",
                            temp_path=getattr(tmp, "name", None),
                            request_id=request_id
                        )
            
            else:
                raise ValueError("File must be a .pptx or .ppt file")
            
            logger.log_file_processing(
                filename=filename,
                file_size=len(file_bytes),
                file_type="powerpoint",
                operation="extraction",
                success=True,
                slides_extracted=len(result),
                request_id=request_id
            )
            
            return result
            
        except Exception as e:
            logger.error(
                "PowerPoint extraction failed",
                filename=filename,
                error=str(e),
                request_id=request_id
            )
            raise

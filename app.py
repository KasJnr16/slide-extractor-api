import pytesseract
from pdf2image import convert_from_bytes
from PyPDF2 import PdfReader
from PIL import Image
import io
import re
import os
import subprocess
from docx import Document
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from fastapi.responses import StreamingResponse
from docx.shared import Pt
import zipfile
from io import BytesIO

app = FastAPI()

# Allow your web app to access this API
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # restrict in production
    allow_methods=["*"],
    allow_headers=["*"],
)


# ======================================================
# HELPER FUNCTIONS (OCR + PDF + DOCX + TXT)
# ======================================================
def clean_text(text):
    text = re.sub(r'\n\s*\n', '\n\n', text)
    text = text.replace("\t", " ")
    return text.strip()


def extract_pdf_text(pdf_bytes):
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text.strip()
    except:
        return ""


def ocr_image(image: Image.Image):
    gray = image.convert("L")
    return pytesseract.image_to_string(gray, lang="eng")


def extract_docx_text(file_bytes):
    try:
        file_stream = io.BytesIO(file_bytes)
        doc = Document(file_stream)
        full_text = [p.text for p in doc.paragraphs]
        return clean_text("\n".join(full_text))
    except Exception as e:
        raise ValueError(f"Error reading DOCX: {e}")


def extract_text_file(file_bytes):
    try:
        return clean_text(file_bytes.decode("utf-8"))
    except:
        return clean_text(file_bytes.decode("latin-1"))


def extract_text_from_any(file_bytes: bytes, filename: str):
    filename = filename.lower()

    # ----- PDF -----
    if filename.endswith(".pdf"):
        extracted = extract_pdf_text(file_bytes)
        if extracted and len(extracted) > 20:
            return clean_text(extracted)

        images = convert_from_bytes(file_bytes)
        ocr_text = ""
        for img in images:
            ocr_text += ocr_image(img) + "\n"
        return clean_text(ocr_text)

    # ----- Images -----
    if filename.endswith((".jpg", ".jpeg", ".png", ".tiff")):
        img = Image.open(io.BytesIO(file_bytes))
        return clean_text(ocr_image(img))

    # ----- DOCX -----
    if filename.endswith(".docx"):
        return extract_docx_text(file_bytes)

    # ----- TXT -----
    if filename.endswith(".txt"):
        return extract_text_file(file_bytes)

    raise ValueError("Unsupported file type for extraction.")


# ------------------------------
# Python .pptx extraction
# ------------------------------
def extract_from_shape(shape, collected):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            collected.append(paragraph.text)

    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    collected.append(paragraph.text)

    if shape.shape_type == 6:  # GROUPED SHAPE
        for subshape in shape.shapes:
            extract_from_shape(subshape, collected)


def extract_text_from_pptx_file(file_like):
    prs = Presentation(file_like)
    all_text = []

    for slide_index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            extract_from_shape(shape, slide_text)

        combined = "\n".join(filter(None, slide_text))
        all_text.append({"slide": slide_index + 1, "text": combined})
    
    return all_text


# ------------------------------
# Java .ppt extraction
# ------------------------------
def extract_text_from_ppt_file(file_path):
    jar_path = os.path.join(os.path.dirname(__file__), "ppt_converter/target/ppt-converter-1.0-jar-with-dependencies.jar")

    
    try:
        result = subprocess.run(
            ["java", "-jar", jar_path, file_path],
            capture_output=True,
            text=True,
            check=True
        )
    except subprocess.CalledProcessError as e:
        print("❌ Java extraction failed:", e)
        print("stdout:", e.stdout)
        print("stderr:", e.stderr)
        return []

    slides_text = []
    current_slide = 0
    slide_lines = []

    for line in result.stdout.splitlines():
        if line.startswith("--- Slide"):
            if slide_lines:
                slides_text.append({"slide": current_slide, "text": "\n".join(slide_lines)})
                slide_lines = []

            current_slide = int(line.replace("--- Slide", "").replace("---", "").strip())
        else:
            slide_lines.append(line)

    if slide_lines:
        slides_text.append({"slide": current_slide, "text": "\n".join(slide_lines)})

    return slides_text


# ----------------------------
# Helper to create zip in memory
# ----------------------------
def create_exam_package_in_memory(document_name, questions, answers):
    """
    Creates questions and answers Word docs in memory and zips them.

    Returns:
        BytesIO: The zip file as a BytesIO object.
    """
    # In-memory Word docs
    questions_io = BytesIO()
    answers_io = BytesIO()

    # ------------------
    # Questions doc
    # ------------------
    doc_q = Document()
    doc_q.add_heading(f"{document_name} - Questions", level=1)
    for q in questions:  # ✅ Removed enumerate
        if q.strip():  # Only add non-empty lines
            para = doc_q.add_paragraph(q)  # ✅ No numbering added here
            para.paragraph_format.space_after = Pt(6)
    doc_q.save(questions_io)
    questions_io.seek(0)  # reset pointer

    # ------------------
    # Answers doc
    # ------------------
    doc_a = Document()
    doc_a.add_heading(f"{document_name} - Answers", level=1)
    for a in answers:  # ✅ Removed enumerate
        if a.strip():  # Only add non-empty lines
            para = doc_a.add_paragraph(a)  # ✅ No numbering added here
            para.paragraph_format.space_after = Pt(6)
    doc_a.save(answers_io)
    answers_io.seek(0)

    # ------------------
    # Create zip in memory
    # ------------------
    zip_io = BytesIO()
    with zipfile.ZipFile(zip_io, mode="w") as zipf:
        zipf.writestr(f"{document_name}_questions.docx", questions_io.getvalue())
        zipf.writestr(f"{document_name}_answers.docx", answers_io.getvalue())
    zip_io.seek(0)

    return zip_io


# ------------------------------
# API endpoint
# ------------------------------
@app.post("/extract_document_text/")
async def extract_document_text(file: UploadFile = File(...)):
    contents = await file.read()
    filename = file.filename.lower()

    try:
        # -------------------------------
        # Handle PPTX
        # -------------------------------
        if filename.endswith(".pptx"):
            file_like = BytesIO(contents)
            slides = extract_text_from_pptx_file(file_like)

            combined = "\n\n".join(
                f"Slide {s['slide']}:\n{s['text']}" for s in slides
            )
            return {
                "filename": filename,
                "text": combined.strip()
            }

        # -------------------------------
        # Handle PPT
        # -------------------------------
        if filename.endswith(".ppt"):
            # Save temp file for Java extractor
            temp_path = os.path.join(os.path.dirname(__file__), filename)
            with open(temp_path, "wb") as f:
                f.write(contents)

            slides = extract_text_from_ppt_file(temp_path)
            os.remove(temp_path)  # Clean temp file

            combined = "\n\n".join(
                f"Slide {s['slide']}:\n{s['text']}" for s in slides
            )
            return {
                "filename": filename,
                "text": combined.strip()
            }

        # -------------------------------
        # Other file types (PDF, DOCX, IMG, TXT)
        # -------------------------------
        extracted = extract_text_from_any(contents, filename)
        return {"filename": filename, "text": extracted}

    except Exception as e:
        return {"error": str(e)}


@app.post("/extract_text/")
async def extract_text(file: UploadFile = File(...)):
    # Ensure only .pptx or .ppt
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in (".pptx", ".ppt"):
        return {"error": "File must be a .pptx or .ppt"}

    # Save uploaded file temporarily if it's a .ppt (Java needs a file path)
    if ext == ".ppt":
        temp_path = os.path.join(os.path.dirname(__file__), file.filename)
        with open(temp_path, "wb") as f:
            f.write(await file.read())

        slides_text = extract_text_from_ppt_file(temp_path)

        # Optionally remove temp file
        os.remove(temp_path)

        return {"filename": file.filename, "slides": slides_text}

    # For .pptx, read in memory
    contents = await file.read()
    file_like = BytesIO(contents)
    slides_text = extract_text_from_pptx_file(file_like)
    
    return {"filename": file.filename, "slides": slides_text}


@app.post("/generate_exam_zip/")
async def generate_exam_zip(
    document_name: str,
    questions_file: UploadFile = File(...),
    answers_file: UploadFile = File(...)
):
    """
    Accepts two text files (questions and answers) and returns a zip
    containing two Word documents (questions + answers) without saving them on disk.
    """
    # Read uploaded text files
    questions_text = (await questions_file.read()).decode("utf-8").splitlines()
    answers_text = (await answers_file.read()).decode("utf-8").splitlines()

    # Create zip in memory
    zip_io = create_exam_package_in_memory(document_name, questions_text, answers_text)

    # Return as downloadable file
    return StreamingResponse(
        zip_io,
        media_type="application/x-zip-compressed",
        headers={"Content-Disposition": f"attachment; filename={document_name}.zip"}
    )
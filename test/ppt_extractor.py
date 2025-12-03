from pptx import Presentation
import tkinter as tk
from tkinter import filedialog
import glob
import os
import subprocess

# ------------------------------------------
# Python extractor for .pptx
# ------------------------------------------
def extract_from_shape(shape, collected):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            collected.append(paragraph.text)

    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    collected.append(paragraph.text)

    if shape.shape_type == 6:  # GROUPED SHAPE
        for subshape in shape.shapes:
            extract_from_shape(subshape, collected)


def extract_text_pptx(filename):
    prs = Presentation(filename)
    all_text = []

    for slide_index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            extract_from_shape(shape, slide_text)

        combined = "\n".join(filter(None, slide_text))
        all_text.append((slide_index + 1, combined))
    
    return all_text

# ------------------------------------------
# Java extractor for .ppt
# ------------------------------------------
def extract_text_ppt(filename):
    jar_path = os.path.join(os.path.dirname(__file__), "../ppt_converter/target/ppt-converter-1.0-jar-with-dependencies.jar")

    try:
        result = subprocess.run(
            ["java", "-jar", jar_path, filename],
            capture_output=True,
            text=True,
            check=True
        )
    except subprocess.CalledProcessError as e:
        print("❌ Java extraction failed:", e)
        print("stdout:", e.stdout)
        print("stderr:", e.stderr)
        return []

    slides_text = []
    current_slide = 0
    slide_lines = []

    for line in result.stdout.splitlines():
        if line.startswith("--- Slide"):
            if slide_lines:
                slides_text.append((current_slide, "\n".join(slide_lines)))
                slide_lines = []

            current_slide = int(line.replace("--- Slide", "").replace("---", "").strip())
        else:
            slide_lines.append(line)

    if slide_lines:
        slides_text.append((current_slide, "\n".join(slide_lines)))

    return slides_text

# ------------------------------------------
# File picker or desktop fallback
# ------------------------------------------
def pick_file_or_desktop():
    root = tk.Tk()
    root.withdraw()

    file_path = filedialog.askopenfilename(
        title="Select PowerPoint File",
        filetypes=[("PowerPoint files", ("*.ppt", "*.pptx"))]
    )

    if file_path:
        print(f"Selected: {file_path}")
        return file_path

    print("No file selected. Checking Desktop...")

    desktop = os.path.join(os.path.expanduser("~"), "Desktop")
    files = glob.glob(desktop + "\\*.pptx") + glob.glob(desktop + "\\*.ppt")

    if not files:
        print("❌ No PPT(PPTX) files found on Desktop.")
        return None

    file_path = files[0]
    print(f"Using Desktop file: {file_path}")
    return file_path

# ------------------------------------------
# Save to TXT
# ------------------------------------------
def save_to_txt(text_data, original_file_path):
    base_dir = r"C:\Users\moji\Desktop\Programming\python\slides_lib"
    txt_folder = os.path.join(base_dir, "txt")

    if not os.path.exists(txt_folder):
        os.makedirs(txt_folder)

    base_name = os.path.splitext(os.path.basename(original_file_path))[0]
    output_path = os.path.join(txt_folder, f"{base_name}.txt")

    with open(output_path, "w", encoding="utf-8") as f:
        for slide_num, text in text_data:
            f.write(f"--- Slide {slide_num} ---\n")
            f.write(text)
            f.write("\n\n")

    print(f"✅ Saved extracted text to: {output_path}")

# ------------------------------------------
# MAIN LOGIC
# ------------------------------------------
file_path = pick_file_or_desktop()

if file_path:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".ppt":
        slides_text = extract_text_ppt(file_path)
    else:
        slides_text = extract_text_pptx(file_path)

    save_to_txt(slides_text, file_path)
else:
    print("\n❌ Could not load any PowerPoint file.")
